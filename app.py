from flask import Flask, render_template, request, redirect, url_for, session, flash, send_from_directory
import sqlite3
import os
import time
from werkzeug.utils import secure_filename
import os
import pypdf
import pytesseract
from docx import Document
from pptx import Presentation   # مكتبة لقراءة PowerPoint
from PIL import Image
import fitz   # مكتبة PyMuPDF لقراءة PDF كصور عند الحاجة
from flask import Flask

app = Flask(__name__)
app.secret_key = 'musaid_secret_key' 

# -----------------------------
# فلتر تنظيف النصوص
# -----------------------------
def clean_text(text):
    import re
    # إبقاء العربية + الإنجليزية + الأرقام + بعض الرموز الأساسية
    text = re.sub(r'[^\u0600-\u06FFa-zA-Z0-9\s.,!?؟]', '', text)
    # إزالة التكرار المبالغ فيه للأحرف
    text = re.sub(r'(.)\1{2,}', r'\1', text)
    return text.strip()

# -----------------------------
# دالة استخراج النص من الملفات
# -----------------------------
def extract_text_from_file(file_path):
    text = ""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            try:
                # محاولة القراءة كنص
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            except Exception:
                # إذا فشلت القراءة، نستخدم OCR
                doc = fitz.open(file_path)
                for page in doc:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # يدعم العربية والإنجليزية
                    text += pytesseract.image_to_string(img, lang="ara+eng")

        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        else:
            text = "⚠️ نوع الملف غير مدعوم حالياً."

    except Exception as e:
        text = f"⚠️ الملف غير صالح للقراءة: {e}"

    return clean_text(text)

# -----------------------------
# دالة التلخيص
# -----------------------------
def summarize_handout(file_name):
    full_path = os.path.join(app.config['UPLOAD_FOLDER'], file_name)
    content = extract_text_from_file(full_path)

    if not content or len(content.strip()) < 20:
        return "لا يوجد محتوى كافي للتلخيص."

    # تلخيص بسيط: أول 3 جمل
    import re
    sentences = re.split(r'[.!?؟]', content)
    summary = " ".join(sentences[:3])

    return summary.strip()

# إعدادات المجلد وقاعدة البيانات
DATABASE = 'musaid_ist.db'
# تم تعديل المسار ليكون مطلقاً لضمان عدم حدوث تضارب في المجلدات
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

import os
import sqlite3
import psycopg2
import psycopg2.extras # استدعاء المكتبة كاملة يحل مشكلة التنبيه في VS Code

import psycopg2.extras # تأكد من وجود هذا السطر في الأعلى

def get_db_connection():
    url = os.environ.get('DATABASE_URL')
    # تصحيح بداية الرابط ليتوافق مع psycopg2
    if url and url.startswith("postgresql://"):
        url = url.replace("postgresql://", "postgres://", 1)
        
    conn = psycopg2.connect(url, sslmode='require')
    # هذا السطر هو "السر" لكي يعمل الكود القديم مع قاعدة البيانات الجديدة
    conn.cursor_factory = psycopg2.extras.RealDictCursor 
    return conn
# --- التعديل الجوهري: حل مشكلة المعاينة والتحميل (PDF) ---
@app.route('/download/<filename>')
def uploaded_file(filename):
    # as_attachment=False تضمن فتح الملف في المتصفح (المعاينة) بدلاً من إجبار التحميل
    # Flask سيتعرف تلقائياً على امتداد .pdf ويرسل الرأس الصحيح للمتصفح
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=False)

# --- 1. واجهة الطالب ---
@app.route('/')
def index():
    conn = get_db_connection()
    depts = conn.execute('SELECT * FROM departments').fetchall()
    conn.close()
    return render_template('index.html', depts=depts)
    
@app.route('/search')
def search():
    dept_id = request.args.get('dept')
    semester = request.args.get('semester')
    
    conn = get_db_connection()
    query = '''
        SELECT h.*, s.subject_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN teachers t ON h.teacher_id = t.id
        WHERE h.dept_id = ? AND h.semester = ?
    '''
    results = conn.execute(query, (dept_id, semester)).fetchall()
    
    dept_name_row = conn.execute('SELECT dept_name FROM departments WHERE id = ?', (dept_id,)).fetchone()
    dept_name = dept_name_row['dept_name'] if dept_name_row else "غير معروف"
    
    processed_results = []
    for row in results:
        item = dict(row)
        # التلخيص العميق باستخدام الملف نفسه
        summary = summarize_handout(item.get('file_path', ''))
        item['flash_summary'] = summary
        processed_results.append(item)
    conn.close()
    
    return render_template('results.html', results=processed_results, dept_name=dept_name, semester=semester)

# --- 2. نظام تسجيل الدخول ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form['email'].strip()
        password = request.form['password'].strip()
        
        conn = get_db_connection()
        user = conn.execute(
            'SELECT * FROM teachers WHERE email = ? AND password = ?', 
            (email, password)
        ).fetchone()
        conn.close()

        if user:
            # حفظ بيانات المستخدم في الجلسة
            session['user_id'] = user['id']
            session['user_name'] = user['full_name']

            # تحديد الدور حسب البريد أو العمود role في قاعدة البيانات
            if email == 'admin@musaid.edu.ly':
                session['role'] = 'admin'
                return redirect(url_for('admin_dashboard'))
            else:
                session['role'] = 'teacher'
                return redirect(url_for('teacher_dashboard'))
        else:
            flash('خطأ في البريد الإلكتروني أو كلمة المرور')
            
    return render_template('login.html')

# --- 3. لوحة تحكم الأستاذ ---
@app.route('/teacher')
def teacher_dashboard():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    depts = conn.execute('SELECT id, dept_name FROM departments').fetchall()
    depts_list = [dict(row) for row in depts]
    
    my_handouts = conn.execute('''
        SELECT h.*, s.subject_name 
        FROM handouts h 
        JOIN subjects s ON h.subject_id = s.id 
        WHERE h.teacher_id = ? 
        ORDER BY h.id DESC
    ''', (session['user_id'],)).fetchall()
    my_handouts_list = [dict(row) for row in my_handouts]

    # عدد المذكرات التي رفعها هذا الأستاذ
    my_count = len(my_handouts_list)

    # أكثر مادة رفع لها مذكرات
    top_subject = None
    if my_count > 0:
        subject_counts = {}
        for h in my_handouts_list:
            subject_counts[h['subject_name']] = subject_counts.get(h['subject_name'], 0) + 1
        top_subject = max(subject_counts, key=subject_counts.get)

    # إجمالي المذكرات لجميع الأساتذة
    total_count = conn.execute("SELECT COUNT(*) FROM handouts").fetchone()[0]

    # نسبة المشاركة
    participation = round((my_count / total_count) * 100, 1) if total_count > 0 else 0

    # تنبيه ذكي عند قلة النشاط (شرط عام فقط)
    ai_message = None
    if my_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم تقم برفع أي مذكرة حتى الآن هذا الفصل."

    conn.close()
    
    return render_template('teacher_dashboard.html', 
                           name=session['user_name'], 
                           depts=depts_list,
                           my_handouts=my_handouts_list,
                           my_count=my_count,
                           top_subject=top_subject,
                           participation=participation,
                           ai_alert=ai_message)

import os
from werkzeug.utils import secure_filename

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))

    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    subject_id = request.form.get('subject_id')
    title = request.form.get('title')
    notes = request.form.get('notes')
    files = request.files.getlist('files[]')

    conn = get_db_connection()
    try:
        # 🔍 فحص التكرار قبل الحفظ
        duplicate = conn.execute("""
            SELECT COUNT(*) FROM handouts 
            WHERE subject_id = ? AND LOWER(title) = LOWER(?)
        """, (subject_id, title)).fetchone()[0]

        if duplicate > 0:
            conn.close()
            flash("⚠️ تنبيه ذكي: هناك مذكرة مشابهة مرفوعة مسبقًا لهذه المادة.")
            return redirect(url_for('teacher_dashboard'))

        # إذا لا يوجد تكرار → نحفظ الملفات
        for file in files:
            if file and file.filename != '':
                original_filename = file.filename
                ext = os.path.splitext(original_filename)[1]
                base_name = secure_filename(os.path.splitext(original_filename)[0])
                unique_filename = f"{int(time.time())}_{base_name}{ext}"
                file.save(os.path.join(app.config['UPLOAD_FOLDER'], unique_filename))

                conn.execute('''
                    INSERT INTO handouts (teacher_id, subject_id, dept_id, semester, title, notes, file_path)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                ''', (session['user_id'], subject_id, dept_id, semester, title, notes, unique_filename))
        
        conn.commit()
        flash('✅ تم رفع الملفات بنجاح!')
    except Exception as e:
        print(f"Error: {e}")
        flash('❌ حدث خطأ أثناء الرفع!')
    finally:
        conn.close()
    
    return redirect(url_for('teacher_dashboard'))

# --- 4. لوحة تحكم المدير ---
@app.route('/admin')
def admin_dashboard():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    teachers = conn.execute("SELECT * FROM teachers WHERE email != 'admin@musaid.edu.ly' ORDER BY id DESC").fetchall()
    conn.close()
    return render_template('admin/dashboard.html', teachers=teachers)

@app.route('/admin/add_teacher', methods=['POST'])
def add_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    full_name = request.form.get('name')
    email = request.form.get('email')
    password = request.form.get('password')
    
    conn = get_db_connection()
    try:
        conn.execute("INSERT INTO teachers (full_name, email, password) VALUES (?, ?, ?)",
                   (full_name, email, password))
        conn.commit()
        flash('تم تسجيل الأستاذ بنجاح!')
    except Exception as e:
        flash(f'خطأ: قد يكون البريد مسجل مسبقاً!')
    finally:
        conn.close()
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/edit_teacher', methods=['POST'])
def edit_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    t_id = request.form.get('id')
    name = request.form.get('name')
    email = request.form.get('email')
    password = request.form.get('password')
    
    conn = get_db_connection()
    conn.execute("UPDATE teachers SET full_name = ?, email = ?, password = ? WHERE id = ?", 
                 (name, email, password, t_id))
    conn.commit()
    conn.close()
    flash('تم تحديث بيانات الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/subjects')
def admin_subjects():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    depts = conn.execute("SELECT * FROM departments WHERE id IN (1, 2)").fetchall()
    query = '''
        SELECT s.id, s.subject_name, cs.dept_id, cs.semester 
        FROM subjects s
        JOIN course_structure cs ON s.id = cs.subject_id
    '''
    subjects = conn.execute(query).fetchall()
    conn.close()
    
    # منطق التنبيه الذكي
    ai_message = None
    if len(subjects) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مواد مسجلة حتى الآن."
    else:
        # مثال: إذا عدد المواد أكبر من عدد الأقسام × 10، نعتبره مؤشر ضغط
        if len(subjects) > len(depts) * 10:
            ai_message = "🔔 تنبيه ذكي: هناك عدد كبير من المواد مقارنة بالأقسام، تحقق من التوزيع."
    
    return render_template('admin/subjects.html',
                           depts=depts,
                           subjects=subjects,
                           ai_alert=ai_message)

@app.route('/admin/delete_teacher/<int:id>')
def delete_teacher(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    conn.execute("DELETE FROM teachers WHERE id = ?", (id,))
    conn.commit()
    conn.close()
    flash('تم حذف حساب الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/add_subject', methods=['POST'])
def add_subject():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
        
    name = request.form.get('subject_name')
    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    
    conn = get_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("INSERT INTO subjects (subject_name) VALUES (?)", (name,))
        subject_id = cursor.lastrowid
        
        conn.execute('''
            INSERT INTO course_structure (subject_id, dept_id, semester) 
            VALUES (?, ?, ?)
        ''', (subject_id, dept_id, semester))
        
        conn.commit()
        flash('تمت إضافة المادة وربطها بالقسم بنجاح!')
    except Exception as e:
        flash('حدث خطأ أثناء إضافة المادة.')
    finally:
        conn.close()
    return redirect(url_for('admin_subjects'))

@app.route('/admin/monitor')
def admin_monitor():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    query = '''
        SELECT h.*, s.subject_name, d.dept_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN departments d ON h.dept_id = d.id
        JOIN teachers t ON h.teacher_id = t.id
        ORDER BY h.id DESC
    '''
    logs = conn.execute(query).fetchall()
    conn.close()
    
    # منطق التنبيه الذكي
    ai_message = None
    if len(logs) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مذكرات مرفوعة حتى الآن."
    elif len(logs) > 20:
        ai_message = "🔔 تنبيه ذكي: تم رفع عدد كبير من المذكرات مؤخرًا، تحقق من صحتها."

    return render_template('admin/monitor.html', logs=logs, ai_alert=ai_message)


@app.route('/admin/delete_subject/<int:id>')
def delete_subject(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        conn.execute("DELETE FROM course_structure WHERE subject_id = ?", (id,))
        conn.execute("DELETE FROM subjects WHERE id = ?", (id,))
        conn.commit()
        flash('تم حذف المادة بنجاح.')
        # تنبيه ذكي عند الحذف
        ai_message = "🔔 تنبيه ذكي: تم حذف مادة من النظام، تأكد من تحديث الهيكل الدراسي."
    except Exception as e:
        flash('حدث خطأ أثناء الحذف.')
        ai_message = "🔔 تنبيه ذكي: حدث خطأ أثناء محاولة حذف المادة."
    finally:
        conn.close()
    
    # نمرر التنبيه مع إعادة التوجيه
    return redirect(url_for('admin_subjects', ai_alert=ai_message))


@app.route('/admin/reports')
def admin_reports():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        t_count = conn.execute("SELECT COUNT(*) FROM teachers WHERE email != 'admin@musaid.edu.ly'").fetchone()[0]
        s_count = conn.execute('SELECT COUNT(*) FROM subjects').fetchone()[0]
        l_count = conn.execute('SELECT COUNT(*) FROM handouts').fetchone()[0]
        it_count = conn.execute('SELECT COUNT(*) FROM handouts WHERE dept_id = 1').fetchone()[0]
        acc_count = conn.execute('SELECT COUNT(*) FROM handouts WHERE dept_id = 2').fetchone()[0]
    except Exception as e:
        t_count = s_count = l_count = it_count = acc_count = 0
    finally:
        conn.close()
    
    # منطق التنبيه الذكي
    ai_message = None
    if l_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم يتم رفع أي مذكرة هذا الأسبوع."
    elif acc_count < 2:
        ai_message = "🔔 تنبيه ذكي: نشاط قسم الإدارة منخفض، يرجى المتابعة."
    
    return render_template('admin/reports.html', 
                           t_count=t_count, 
                           s_count=s_count, 
                           l_count=l_count,
                           it_count=it_count,
                           acc_count=acc_count,
                           ai_alert=ai_message)

@app.route('/logout')
def logout():
    session.clear() 
    flash('تم تسجيل الخروج بنجاح.')
    return redirect(url_for('login'))

# --- الدوال المساعدة للجافاسكريبت ---
@app.route('/get_semesters/<int:dept_id>')
def get_semesters(dept_id):
    conn = get_db_connection()
    semesters = conn.execute('''
        SELECT DISTINCT semester FROM course_structure 
        WHERE dept_id = ? ORDER BY semester
    ''', (dept_id,)).fetchall()
    conn.close()
    if not semesters:
        return {"semesters": [{"id": i, "number": i} for i in range(1, 7)]}
    return {"semesters": [{"id": s['semester'], "number": s['semester']} for s in semesters]}

@app.route('/get_subjects/<int:dept_id>/<int:semester>')
def get_subjects(dept_id, semester):
    conn = get_db_connection()
    try:
        subjects = conn.execute('''
            SELECT s.id, s.subject_name FROM subjects s
            JOIN course_structure cs ON s.id = cs.subject_id
            WHERE cs.dept_id = ? AND cs.semester = ?
        ''', (dept_id, semester)).fetchall()
    except:
        subjects = []
    conn.close()
    return {"subjects": [{"id": s['id'], "name": s['subject_name']} for s in subjects]}

@app.route('/delete_handout/<int:handout_id>')
def delete_handout(handout_id):
    if 'user_id' not in session or session['role'] not in ['teacher', 'admin']:
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        handout = conn.execute('SELECT * FROM handouts WHERE id = ?', (handout_id,)).fetchone()
        if handout:
            if session['role'] == 'admin' or handout['teacher_id'] == session['user_id']:
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], handout['file_path'])
                if os.path.exists(file_path):
                    os.remove(file_path)
                conn.execute('DELETE FROM handouts WHERE id = ?', (handout_id,))
                conn.commit()
                flash('تم حذف المذكرة بنجاح')
    except Exception as e:
        flash(f'خطأ: {str(e)}')
    finally:
        conn.close()
    return redirect(request.referrer or url_for('index'))

# تجهيز حساب المدير
with app.app_context():
    conn = sqlite3.connect(DATABASE)
    conn.execute("UPDATE teachers SET password = '33557799' WHERE email = 'admin@musaid.edu.ly'")
    conn.execute('''
        INSERT OR IGNORE INTO teachers (full_name, email, password) 
        VALUES ('مدير النظام', 'admin@musaid.edu.ly', '33557799')
    ''')
    conn.commit()
    conn.close()
    if __name__ == "__main__":
    # تأكد أن السطور التالية تبدأ بـ 4 مسافات (وليس Tab)
        port = int(os.environ.get("PORT", 5000))
        # تشغيل التطبيق مع تعطيل الديباج في السيرفر
        app.run(host='0.0.0.0', port=port, debug=False)